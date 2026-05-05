# tests/test_bp_parser.py
import pytest
from unittest.mock import patch, MagicMock
from core.bp_parser import extract_text_from_pdf, extract_text_from_pptx, extract_project_info, extract_report_info

def test_extract_text_from_pdf(tmp_path):
    fake_pdf = tmp_path / "test.pdf"
    fake_pdf.write_bytes(b"")
    result = extract_text_from_pdf(str(fake_pdf))
    assert isinstance(result, str)

def test_extract_text_from_pptx(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "测试项目"
    slide.placeholders[1].text = "我们是一家AI公司，专注于工业质检。"
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))
    text = extract_text_from_pptx(str(pptx_path))
    assert "测试项目" in text
    assert "工业质检" in text

def test_extract_project_info_returns_required_fields(mocker):
    mocker.patch("core.llm.call_llm", return_value='{"name": "智检科技", "sector": "AI", "sub_sector": "AI+工业质检", "stage": "Pre-A", "location": "深圳", "description": "基于AI的工业质检系统", "highlights": "准确率99.5%，已服务50家工厂", "financing_need": "3000万人民币"}')
    result = extract_project_info("这是一段BP文本内容")
    assert result["sector"] == "AI"
    assert result["sub_sector"] == "AI+工业质检"
    assert result["stage"] == "Pre-A"
    assert "name" in result
    assert "description" in result

def test_extract_project_info_handles_partial_extraction(mocker):
    mocker.patch("core.llm.call_llm", return_value='{"name": "某项目", "sector": "AI"}')
    result = extract_project_info("短文本")
    assert result.get("sub_sector", "") == ""
    assert result["sector"] == "AI"


def test_extract_report_info_returns_six_fields(mocker):
    fake_response = '{"founded_year":"2018年","headquarters":"北京","sector":"先进制造","main_products":"工业质检","team":"张三，CEO","customers":"富士康"}'
    mocker.patch("core.llm.call_llm", return_value=fake_response)
    result = extract_report_info("some bp text")
    assert result["founded_year"] == "2018年"
    assert result["headquarters"] == "北京"
    assert result["sector"] == "先进制造"
    assert result["main_products"] == "工业质检"
    assert result["team"] == "张三，CEO"
    assert result["customers"] == "富士康"


def test_extract_report_info_defaults_missing_fields(mocker):
    fake_response = '{"founded_year":"2020年"}'
    mocker.patch("core.llm.call_llm", return_value=fake_response)
    result = extract_report_info("text")
    assert result["headquarters"] == ""
    assert result["sector"] == ""
    assert result["main_products"] == ""
    assert result["team"] == ""
    assert result["customers"] == ""


def test_extract_report_info_handles_json_in_markdown(mocker):
    fake_response = '```json\n{"founded_year":"2019年","headquarters":"上海","sector":"AI","main_products":"x","team":"y","customers":"z"}\n```'
    mocker.patch("core.llm.call_llm", return_value=fake_response)
    result = extract_report_info("text")
    assert result["founded_year"] == "2019年"


def test_extract_report_info_handles_bad_json(mocker):
    mocker.patch("core.llm.call_llm", return_value="not json at all")
    result = extract_report_info("text")
    assert result == {"founded_year": "", "headquarters": "", "sector": "",
                      "main_products": "", "team": "", "customers": ""}
