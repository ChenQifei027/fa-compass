# core/bp_parser.py
import json
import re


def _parse_json_safe(raw: str) -> dict:
    # 去掉 markdown 代码围栏，定位 { ... }
    text = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw.strip(), flags=re.MULTILINE).strip()
    start, end = text.find("{"), text.rfind("}")
    if start == -1 or end == -1:
        return {}
    text = text[start: end + 1]
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # 修复尾部逗号
    text = re.sub(r",\s*([}\]])", r"\1", text)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        return {}

def extract_text_from_pdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append(f"[第{idx}页]")
                texts.extend(slide_texts)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_from_file(file_path: str) -> str:
    path = file_path.lower()
    if path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif path.endswith((".pptx", ".ppt")):
        return extract_text_from_pptx(file_path)
    return ""

EXTRACT_PROMPT = """你是一名专业的投资 FA 助手。请从以下 BP（商业计划书）文本中提取关键信息，以 JSON 格式返回。

提取规则：
- name: 公司/项目全称（不含"股份"、"有限公司"等后缀），通常在封面或首页
- sector: 一级赛道，用简短标签（如：AI、消费、医疗健康、企业服务、硬件、新能源、先进制造等）
- sub_sector: 细分方向（比 sector 更具体，如：AI+工业质检、智慧交通、分布式数据库、新能源储能等）
- stage: 当前融资阶段，优先提取"本轮融资""当前轮次"等明确表述，枚举为：天使轮/Pre-A轮/A轮/A+轮/Pre-B轮/B轮/B+轮/C轮/C+轮/D轮及以后/上市前/战略轮
- location: 公司注册地或总部所在城市（只写城市名，如"北京"、"上海"、"成都"）
- description: 用一句话描述公司核心业务（≤80字，第三人称）
- highlights: 3-5个核心亮点，用分号"；"分隔（每条≤30字，优先写营收、客户、专利、资质等可量化信息）
- financing_need: 本轮融资金额，格式如"3000万元"或"1亿元"；如未明确说明则返回空字符串

字段找不到时返回空字符串，不要猜测。只返回合法 JSON，不要其他内容。
注意：JSON 字段值中不能出现裸 ASCII 双引号，引用专有名词时用书名号《》或直接省略引号。

BP 文本：
{text}"""


def extract_project_info(text: str, api_key: str = "") -> dict:
    from core.llm import call_llm
    raw = call_llm(EXTRACT_PROMPT.format(text=text[:15000]))
    data = _parse_json_safe(raw)
    defaults = {"name": "", "sector": "", "sub_sector": "", "stage": "",
                "location": "", "description": "", "highlights": "", "financing_need": ""}
    return {**defaults, **data}


REPORT_PROMPT = """你是一名专业的投资 FA 助手，正在为项目撰写尽调报告。从以下 BP 文本中提取项目基本信息，以 JSON 返回。

提取规则：
- founded_year: 公司成立年份，4位数字（如"2018"），找不到返回空字符串
- headquarters: 总部或注册地城市（如"北京"、"上海·浦东新区"），找不到返回空字符串
- sector: 细化赛道描述（比一级赛道更具体，如"分布式实时分析型AI数据库"、"智慧交通+计算机视觉"）
- main_products: 核心产品或服务名称及简述（≤150字，如有多个产品用"；"分隔）
- team: 核心创始团队成员姓名、职务、关键背景（≤200字，如"CEO张某，前阿里P9；CTO李某，清华博士"）
- customers: 已签约或在谈的主要客户/合作方（≤100字，机构名用顿号分隔）

字段找不到时返回空字符串，不要编造。只返回合法 JSON，不要其他内容。
注意：JSON 字段值中不能出现裸 ASCII 双引号，如需引用专有名词请用书名号《》。

BP 文本：
{text}"""


def extract_report_info(text: str) -> dict:
    from core.llm import call_llm
    raw = call_llm(REPORT_PROMPT.format(text=text[:15000]))
    data = _parse_json_safe(raw)
    defaults = {"founded_year": "", "headquarters": "", "sector": "",
                "main_products": "", "team": "", "customers": ""}
    return {**defaults, **{k: v for k, v in data.items() if k in defaults}}
